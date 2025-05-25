from pathlib import Path
from django.http import HttpResponse
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
import tempfile
import logging
import shutil
import os
import pandas as pd
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader

from services.openai_client import OpenAIClient
from services.groq_client import GroqApiClient
from services.podman_executor import PodmanExecutor

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class RunProgramView(APIView):
    # Class attribute for PodmanExecutor (Singleton pattern)
    podman_executor = None

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        # Initialize PodmanExecutor once for the whole class
        if RunProgramView.podman_executor is None:
            logger.info("Initializing PodmanExecutor for the first time.")
            RunProgramView.podman_executor = PodmanExecutor()
        else:
            logger.info("Reusing existing PodmanExecutor instance.")

    def post(self, request, *args, **kwargs):
        logger.info("Starting the post request for RunProgramView.")
        
        # Create a temporary directory for the upload process
        logger.info("Creating a temporary directory.")
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_directory = Path(temp_dir)
            logger.info(f"Temporary directory created at: {temp_directory}")

            # Create an output directory within the temporary directory
            output_directory = temp_directory / "output"
            output_directory.mkdir(parents=True, exist_ok=True)  # Create output directory
            logger.info(f"Output directory created at: {output_directory}")

            # Handle file uploads
            uploaded_files = self.handle_file_uploads(request.FILES.getlist('files'), temp_directory)
            if not uploaded_files:
                logger.error("No files uploaded.")
                return Response({"error": "No files uploaded."}, status=status.HTTP_400_BAD_REQUEST)

            logger.info(f"Uploaded files: {uploaded_files}")

            # Generate the input files description based on the uploaded files
            input_files_description = self.generate_input_files_description(uploaded_files, temp_directory, 16)
            
            # Check if there's any description generated
            if not input_files_description:
                logger.error("Failed to generate file descriptions.")
                return Response({"error": "Failed to generate file descriptions."}, status=status.HTTP_400_BAD_REQUEST)

            # Generate and save the Python code
            instruction = request.data.get("instruction")
            if not instruction:
                logger.error("Instruction is required.")
                return Response({"error": "Instruction is required"}, status=status.HTTP_400_BAD_REQUEST)
            
            # openai_client = OpenAIClient()
            # generated_code = openai_client.generate_python_code(input_files_description, instruction)

            groq_client = GroqApiClient()
            generated_code = groq_client.generate_python_code(input_files_description, instruction)

            code_file_path = self.save_generated_code(generated_code, temp_directory)
            logger.info(f"Generated code saved to: {code_file_path}")

            # Execute the generated Python script using the class-level PodmanExecutor
            logger.info("Executing the generated Python script.")
            execution_successfull, logs = RunProgramView.podman_executor.execute_script(temp_directory)

            
            if not execution_successfull:
                print('--- error logs ---\n')
                print(logs)
                print('--- error logs ---\n')
                number_of_generation_retries = 2
                logs = ""
                for _ in range(number_of_generation_retries):
                    # Execute the generated Python script
                    generated_code = groq_client.fix_generated_code(generated_code, logs)
                    code_file_path = self.save_generated_code(generated_code, temp_directory)
                    execution_successfull, logs = RunProgramView.podman_executor.execute_script(temp_directory)
                    if execution_successfull:
                        break
                    else:
                        print('--- error logs ---\n')
                        print(logs)
                        print('--- error logs ---\n')
            

            if not execution_successfull:  # Check if there's an error
                logger.error("Execution of the Python script failed.")
                return Response(logs, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

            logger.info("Execution of the Python script successful.")

            # Create a zip file of the output directory and return it as a response
            logger.info("Creating zip response of the output directory.")
            return self.create_download_response(output_directory)  # Pass the output directory

    def generate_input_files_description(self, uploaded_files, temp_directory, num_lines):
        """
        Generate a description string of the uploaded files, including their first few lines.
        This handles both file objects (like UploadedFile) and file paths (str).
        """
        description = ""

        for file in uploaded_files:
            file_name = os.path.basename(file)
            file_extension = os.path.splitext(file_name)[1].lower()
            file_path = os.path.join(temp_directory, file_name)

            description += f"{file_name}:\n\"\"\"\n"
            first_lines = ""

            try:
                # Handle different file types
                if file_extension in ['.txt', '']:  # Text file or unknown extension defaults to text
                    first_lines = self.read_text_file(open(file_path, 'rb'), num_lines)
                elif file_extension == '.docx':  # Word file
                    first_lines = self.read_word_file(open(file_path, 'rb'), num_lines)
                elif file_extension == '.xlsx':  # Excel file
                    first_lines = self.read_excel_file(open(file_path, 'rb'), num_lines)
                elif file_extension == '.pptx':  # PowerPoint file
                    first_lines = self.read_ppt_file(open(file_path, 'rb'), num_lines)
                elif file_extension == '.pdf':  # PDF file
                    first_lines = self.read_pdf_file(open(file_path, 'rb'), num_lines)
                else:
                    first_lines = "Cannot read this file."  # Default to text file for unknown types

                description += first_lines
                description += "...\n\"\"\"\n\n"

            except Exception as e:
                logger.error(f"Error reading file {file_name}: {str(e)}")
                description += f"Error reading file {file_name}\n\n"

        return description

    def read_text_file(self, file, num_lines):
        """Read the first few lines of a text file."""
        first_lines = ""
        for i, line in enumerate(file):
            first_lines += line.decode('utf-8')  # Assuming file encoding is UTF-8
            if i >= num_lines - 1:  # Read up to the number of lines specified
                break
        return first_lines

    def read_word_file(self, file, num_lines):
        """Read the first few paragraphs of a Word (.docx) file."""
        doc = Document(file)
        first_paragraphs = '\n'.join([para.text for para in doc.paragraphs[:num_lines]])  # First num_lines paragraphs
        return first_paragraphs

    def read_excel_file(self, file, num_lines):
        """Read the first few rows of an Excel (.xlsx) file."""
        wb = load_workbook(file, read_only=True)
        sheet = wb.active
        first_rows = '\n'.join([str([cell.value for cell in row]) for row in sheet.iter_rows(min_row=1, max_row=num_lines)])  # First num_lines rows
        return first_rows

    def read_ppt_file(self, file, num_lines):
        """Read the first few slides of a PowerPoint (.pptx) file."""
        prs = Presentation(file)
        first_slides = ""
        for i, slide in enumerate(prs.slides):
            slide_text = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            first_slides += slide_text + '\n'
            if i >= num_lines - 1:  # Read up to the number of slides specified
                break
        return first_slides

    def read_pdf_file(self, file, num_lines):
        """Read the first few pages of a PDF file."""
        reader = PdfReader(file)
        first_pages = ""
        for i, page in enumerate(reader.pages):
            first_pages += page.extract_text() + '\n'
            if i >= num_lines - 1:  # Read up to the number of pages specified
                break
        return first_pages

    def handle_file_uploads(self, files, temp_directory):
        """Handle file uploads and save them to the temporary directory."""
        uploaded_files = []
        for file in files:
            self.save_uploaded_file(file, temp_directory)
            uploaded_files.append(file.name)  # Save file name directly
        return uploaded_files

    def save_uploaded_file(self, file, temp_directory):
        """Save the uploaded file in the temporary directory."""
        file_path = temp_directory / file.name
        with file_path.open('wb') as f:
            f.write(file.read())

    def save_generated_code(self, generated_code, temp_directory):
        """Save the generated code to a file in the temporary directory."""
        code_file_path = temp_directory / "main.py"
        with code_file_path.open('w') as f:
            f.write(generated_code)
        return code_file_path
    
    def create_download_response(self, output_directory):
        """Create a zip file of the output directory and return it as a response."""
        zip_filename = f"output.zip"
        zip_file_path = output_directory.parent / zip_filename  # Store zip file at the temp directory level

        # Get the list of files in the output directory
        files = os.listdir(output_directory)
        
        # Check if there is exactly one file in the directory
        if len(files) == 1:
            # Return the single file directly
            single_file_path = output_directory / files[0]
            with open(single_file_path, 'rb') as single_file:
                response = HttpResponse(single_file.read(), content_type='application/octet-stream')
                response['Content-Disposition'] = f'attachment; filename="{files[0]}"'
                logger.info(f"Returning single file: {files[0]}")
            return response
        elif len(files) > 1:
            # Create a zip file from the output directory
            shutil.make_archive(zip_file_path.with_suffix(''), 'zip', str(output_directory))
            logger.info(f"Created zip file at: {zip_file_path}")

            # Prepare the zip file for download
            with open(zip_file_path, 'rb') as zip_file:
                response = HttpResponse(zip_file.read(), content_type='application/zip')
                response['Content-Disposition'] = f'attachment; filename="{zip_filename}"'
                logger.info(f"Returning zip file: {zip_filename}")
            return response
        else:
            # Optionally handle the case where there are no files
            logger.warning("No files found in the output directory.")
            return HttpResponse("No files available for download.", status=404)
